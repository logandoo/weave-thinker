# Copyright (c) 2026 Weave Thinker Contributors
# SPDX-License-Identifier: Apache-2.0

"""Parse uploaded files (docx, pptx, xlsx, csv, pdf, md) to Markdown."""
import asyncio
import base64
import csv
import io
import logging
import mimetypes
import os
import tempfile
from pathlib import Path

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".csv", ".pdf"}
IMAGE_EXTENSIONS = {".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".svg", ".tiff", ".ico"}


def detect_file_type(filename: str) -> str | None:
    ext = Path(filename).suffix.lower()
    if ext in IMAGE_EXTENSIONS:
        return "image"
    if ext in SUPPORTED_EXTENSIONS:
        return ext.lstrip(".")
    return None


def is_image_pdf(file_path: str) -> bool:
    try:
        from pdfminer.pdfpage import PDFPage
        from pdfminer.pdfparser import PDFParser
        from pdfminer.pdfdocument import PDFDocument
        from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
        from pdfminer.converter import PDFPageAggregator

        with open(file_path, "rb") as f:
            parser = PDFParser(f)
            doc = PDFDocument(parser)
            rsrcmgr = PDFResourceManager()
            device = PDFPageAggregator(rsrcmgr)
            interpreter = PDFPageInterpreter(rsrcmgr, device)

            pages_checked = 0
            total_text_chars = 0
            for page in PDFPage.create_pages(doc):
                interpreter.process_page(page)
                total_text_chars += len(device.get_result().get_text() or "")
                pages_checked += 1
                if pages_checked >= 5:
                    break

            if pages_checked == 0:
                return True
            avg_chars = total_text_chars / pages_checked
            return avg_chars < 10
    except Exception:
        return False


def parse_docx(file_path: str) -> str:
    from docx import Document
    import base64
    import mimetypes

    doc = Document(file_path)
    parts: list[str] = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            parts.append("")
            continue
        style_name = (para.style.name or "").lower() if para.style else ""
        if "heading 1" in style_name:
            parts.append(f"# {text}")
        elif "heading 2" in style_name:
            parts.append(f"## {text}")
        elif "heading 3" in style_name:
            parts.append(f"### {text}")
        elif "heading 4" in style_name:
            parts.append(f"#### {text}")
        elif "heading 5" in style_name:
            parts.append(f"##### {text}")
        elif "heading 6" in style_name:
            parts.append(f"###### {text}")
        elif "list" in style_name:
            parts.append(f"- {text}")
        else:
            parts.append(text)

    for table in doc.tables:
        rows_data: list[list[str]] = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            rows_data.append(cells)
        if rows_data:
            header = rows_data[0]
            parts.append("")
            parts.append("| " + " | ".join(header) + " |")
            parts.append("| " + " | ".join(["---"] * len(header)) + " |")
            for row in rows_data[1:]:
                padded = row + [""] * (len(header) - len(row))
                parts.append("| " + " | ".join(padded[:len(header)]) + " |")

    # Extract inline images and append as base64 markdown.
    try:
        image_markdowns: list[str] = []
        for shape in doc.inline_shapes:
            rId = shape._inline.graphic.graphicData.pic.blipFill.blip.embed
            image_part = doc.part.related_parts.get(rId)
            if not image_part:
                continue
            data = image_part.blob
            ext = os.path.splitext(image_part.partname)[1].lower() or ".png"
            mime = mimetypes.types_map.get(ext, "image/png")
            b64 = base64.b64encode(data).decode("ascii")
            image_markdowns.append(f"![image](data:{mime};base64,{b64})")
        if image_markdowns:
            parts.append("")
            parts.extend(image_markdowns)
    except Exception as e:
        logger.warning("Failed to extract DOCX images: %s", e)

    return "\n".join(parts).strip()


def parse_pptx(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"## 幻灯片 {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                table = shape.table
                rows_data: list[list[str]] = []
                for row in table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    rows_data.append(cells)
                if rows_data:
                    header = rows_data[0]
                    parts.append("")
                    parts.append("| " + " | ".join(header) + " |")
                    parts.append("| " + " | ".join(["---"] * len(header)) + " |")
                    for row in rows_data[1:]:
                        padded = row + [""] * (len(header) - len(row))
                        parts.append("| " + " | ".join(padded[:len(header)]) + " |")
        parts.append("")

    return "\n".join(parts).strip()


def parse_xlsx(file_path: str) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    parts: list[str] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"## {sheet_name}")
        first_row = True
        for row in ws.iter_rows(values_only=True):
            cells = [str(c).strip() if c is not None else "" for c in row]
            if all(c == "" for c in cells):
                continue
            parts.append("| " + " | ".join(cells) + " |")
            if first_row:
                parts.append("| " + " | ".join(["---"] * len(cells)) + " |")
                first_row = False
        parts.append("")

    wb.close()
    return "\n".join(parts).strip()


def parse_csv_file(file_path: str) -> str:
    parts: list[str] = []
    for encoding in ("utf-8", "utf-8-sig", "gbk", "latin-1"):
        try:
            with open(file_path, "r", encoding=encoding, newline="") as f:
                sample = f.read(4096)
                f.seek(0)
                try:
                    dialect = csv.Sniffer().sniff(sample)
                except csv.Error:
                    dialect = csv.excel
                reader = csv.reader(f, dialect)
                first_row = True
                for row in reader:
                    cells = [c.strip() for c in row]
                    parts.append("| " + " | ".join(cells) + " |")
                    if first_row:
                        parts.append("| " + " | ".join(["---"] * len(cells)) + " |")
                        first_row = False
            break
        except (UnicodeDecodeError, UnicodeError):
            continue
    return "\n".join(parts).strip()


def _pdf_tables_to_markdown(tables) -> str:
    out: list[str] = []
    for t in tables or []:
        rows = [[(c or '').replace('\n', ' ').replace('\\', '\\\\').replace('|', '\\|').strip() for c in row] for row in t]
        rows = [r for r in rows if any(r)]
        if not rows:
            continue
        width = max(len(r) for r in rows)
        rows = [r + [''] * (width - len(r)) for r in rows]
        out.append('| ' + ' | '.join(rows[0]) + ' |')
        out.append('|' + '---|' * width)
        for r in rows[1:]:
            out.append('| ' + ' | '.join(r) + ' |')
    return '\n'.join(out)


def parse_pdf(file_path: str) -> str:
    # 许可合规（纯 permissive）：主引擎 = pdfplumber（MIT，内核即 pdfminer.six +
    # 行级排版 + 表格检测），pdfminer 裸文本兜底。禁止在此文件引入 AGPL/GPL PDF 库。
    try:
        import pdfplumber
        parts: list[str] = []
        with pdfplumber.open(file_path) as pdf:
            for idx, page in enumerate(pdf.pages, 1):
                text = page.extract_text(x_tolerance=2.0, y_tolerance=3.0) or ''
                if text.strip():
                    parts.append(text.strip())
                md_tables = _pdf_tables_to_markdown(page.extract_tables())
                if md_tables:
                    parts.append(f"### 表格（第 {idx} 页）\n\n{md_tables}")
        md = '\n\n'.join(parts)
        if md.strip():
            return md
    except Exception as e:
        logger.warning("pdfplumber failed, falling back to pdfminer: %s", e)

    from pdfminer.high_level import extract_text

    text = extract_text(file_path)
    if not text or not text.strip():
        return ""
    lines = text.split("\n")
    cleaned: list[str] = []
    for line in lines:
        stripped = line.strip()
        cleaned.append(stripped if stripped else "")
    return "\n".join(cleaned).strip()


async def parse_file(file_path: str, filename: str) -> dict:
    file_type = detect_file_type(filename)
    if file_type == "image":
        return {
            "success": False,
            "error": "暂不支持解析图片文件，请上传 Word、PPT、Excel、CSV 或 PDF 文档。",
            "file_type": "image",
            "filename": filename,
        }

    if file_type is None:
        ext = Path(filename).suffix.lower()
        return {
            "success": False,
            "error": f"不支持的文件格式: {ext}。支持: .docx, .doc, .pptx, .ppt, .xlsx, .xls, .csv, .pdf",
            "file_type": ext,
            "filename": filename,
        }

    if file_type == "pdf":
        is_image = await asyncio.to_thread(is_image_pdf, file_path)
        if is_image:
            return {
                "success": False,
                "error": "该 PDF 为图片格式，暂不支持解析。请上传文字版 PDF 或其他文档格式。",
                "file_type": "pdf",
                "filename": filename,
            }

    try:
        if file_type == "docx":
            markdown = await asyncio.to_thread(parse_docx, file_path)
        elif file_type == "doc":
            markdown = await asyncio.to_thread(parse_docx, file_path)
        elif file_type == "pptx":
            markdown = await asyncio.to_thread(parse_pptx, file_path)
        elif file_type == "ppt":
            markdown = await asyncio.to_thread(parse_pptx, file_path)
        elif file_type in ("xlsx", "xls"):
            markdown = await asyncio.to_thread(parse_xlsx, file_path)
        elif file_type == "csv":
            markdown = await asyncio.to_thread(parse_csv_file, file_path)
        elif file_type == "pdf":
            markdown = await asyncio.to_thread(parse_pdf, file_path)
        else:
            return {
                "success": False,
                "error": f"暂不支持解析 {file_type} 格式",
                "file_type": file_type,
                "filename": filename,
            }

        if not markdown.strip():
            return {
                "success": False,
                "error": "文件内容为空或无法提取文本",
                "file_type": file_type,
                "filename": filename,
            }

        return {
            "success": True,
            "markdown": markdown,
            "file_type": file_type,
            "filename": filename,
        }
    except Exception as e:
        logger.exception("Failed to parse file %s: %s", filename, e)
        return {
            "success": False,
            "error": f"解析文件失败: {str(e)[:200]}",
            "file_type": file_type,
            "filename": filename,
        }
